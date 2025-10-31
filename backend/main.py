#
# Akıllı İçerik Platformu (Versiyon 4.0 - Kalıcı PostgreSQL Veritabanı)
# Created by b!g
#

# --- Temel Kütüphane İçe Aktarımları ---
import os
import uvicorn
import io
import base64
import json 
import secrets 
from fastapi import FastAPI, UploadFile, File, HTTPException, Header, Depends, status
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.concurrency import run_in_threadpool # Tespit 2.3 (Async Bloklama) Çözümü
from dotenv import load_dotenv 
from typing import Optional, List

# --- Veritabanı İçe Aktarımları (Yeni V2 Mimarisi) ---
from sqlalchemy.orm import Session
from . import crud, models, schemas
from .database import SessionLocal, engine, init_db

# --- Proje Fonksiyonelliği İçin Gerekli İçe Aktarımlar ---
from slugify import slugify
from datetime import datetime

# --- İçerik Okuyucular ve LLM ---
import openai
import PyPDF2         
import docx           
import pptx           
import pytube         

# --- Bulut Depolama Kütüphanesi ---
from google.cloud import storage 

# --- GÜVENLİK VE GCS AYARLARI ---
MAX_FILE_SIZE_MB = 50 
ALLOWED_EXTENSIONS = {
    ".mp3", ".wav", ".m4a", ".pdf", ".docx", ".doc", ".pptx", ".ppt", 
    ".jpg", ".jpeg", ".png"
}
GCS_KEY_ENV_VAR = "GCS_SA_KEY" 
GCS_BUCKET_NAME = "akilli-icerik-raporlari-bbkgzn" # Kendi bucket adınız

# --- API Anahtarını Yükleme ---
load_dotenv(os.path.join(os.path.dirname(__file__), '..', '.env'))
openai.api_key = os.getenv("OPENAI_API_KEY")

if openai.api_key is None and not os.getenv("RENDER"): 
    raise EnvironmentError("OPENAI_API_KEY .env dosyasında ayarlanmamış.")

client = openai.OpenAI(api_key=openai.api_key or os.getenv("OPENAI_API_KEY"))

# --- VERİTABANI BAŞLATMA ---
# NOT: Artık USERS_DB veya load_users() yok.
# Tabloların veritabanında oluşturulmasını sağlıyoruz.
try:
    init_db()
except Exception as e:
    print(f"Veritabanı başlatma hatası (uygulama başlarken): {e}")
    # Render'ın yeniden başlatma döngüsüne girmemesi için devam et
    
# --- FastAPI Sunucusunu Başlatma ---
app = FastAPI(
    title="Akıllı İçerik Platformu API (V2 - Kalıcı DB)",
    description="Çoklu ortam dosyalarını analiz edip kişiselleştirilmiş raporlar oluşturan platform.",
    version="0.4.0" 
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], allow_credentials=True,
    allow_methods=["*"], allow_headers=["*"],
)

# --- VERİTABANI BAĞIMLILIĞI (Dependency Injection) ---
def get_db():
    """
    Her API isteği için bağımsız bir veritabanı oturumu (session) açar
    ve işlem bittiğinde (hata alsa bile) kapatır.
    """
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# --- YENİ GÜVENLİK VE KİMLİK DOĞRULAMA (Tespit 2.2 Çözümü) ---
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token") # /token endpoint'ini kullanır

async def get_current_user_from_token(
    db: Session = Depends(get_db), 
    token: str = Header(..., alias="X-API-TOKEN")
) -> models.User:
    """
    X-API-TOKEN başlığından gelen token'ı alır,
    veritabanında (crud.py) arar ve ilgili kullanıcıyı döndürür.
    Bu, eski get_user_id'nin yerini alır ve çok daha güvenlidir.
    """
    user = crud.get_user_by_token(db, token=token)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Geçersiz veya süresi dolmuş X-API-TOKEN",
            headers={"WWW-Authenticate": "Bearer"},
        )
    return user

# --- AKILLI DOSYA OKUMA İŞLEVLERİ (DEĞİŞİKLİK YOK) ---
# read_audio, read_pdf, read_docx, read_pptx, read_image, download_youtube_audio...
# (Bu fonksiyonlar bir önceki kodla aynı, buraya kopyalamıyorum,
# ancak tam kodda olmalılar. Eğer main.py'yi SİLİP YAPIŞTIRIYORSANIZ,
# bu fonksiyonları ÖNCEKİ main.py'den kopyalayıp buraya ekleyin.)
# --- (OKUMA FONKSİYONLARI BURADA) ---
def read_audio(file_data: UploadFile) -> str:
    try:
        transcription = client.audio.transcriptions.create(
            model="whisper-1", 
            file=(file_data.filename, file_data.file, file_data.content_type)
        )
        return transcription.text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Whisper Okuma Hatası: {e}")

def read_pdf(file_data: UploadFile) -> str:
    full_text = []
    try:
        reader = PyPDF2.PdfReader(file_data.file)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                full_text.append(text)
        return "\n".join(full_text)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF Okuma Hatası: Dosya bozuk veya şifreli olabilir. {e}")

def read_docx(file_data: UploadFile) -> str:
    full_text = []
    try:
        document = docx.Document(file_data.file)
        for para in document.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DOCX Okuma Hatası: {e}")

def read_pptx(file_data: UploadFile) -> str:
    full_text = []
    try:
        presentation = pptx.Presentation(file_data.file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        full_text.append(shape.text)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame.text.strip():
                    full_text.append(notes_slide.notes_text_frame.text)
        return "\n".join(full_text)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPTX Okuma Hatası: {e}")

def read_image(file_data: UploadFile) -> str:
    try:
        image_bytes = file_data.file.read()
        base64_image = base64.b64encode(image_bytes).decode('utf-8')
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "user", "content": [
                    {"type": "text", "text": "Bu görseldeki tüm metni eksiksiz bir şekilde OCR yaparak metin olarak çıkarın. Çıkan metinle ilgili yorum yapmayın, sadece metni döndürün."},
                    {"type": "image_url", "image_url": {"url": f"data:{file_data.content_type};base64,{base64_image}"}}
                ]}], max_tokens=4096)
        return response.choices[0].message.content
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Görsel (OCR) Okuma Hatası: {e}")

def download_youtube_audio(url: str) -> str:
    temp_dir = "./temp"
    os.makedirs(temp_dir, exist_ok=True)
    temp_file_path = None
    try:
        yt = pytube.YouTube(url)
        audio_stream = yt.streams.get_audio_only()
        temp_file_path = audio_stream.download(output_path=temp_dir, filename_prefix="yt_")
        with open(temp_file_path, "rb") as audio_file:
            transcription = client.audio.transcriptions.create(model="whisper-1", file=audio_file)
            return transcription.text
    except pytube.exceptions.VideoUnavailable:
        raise HTTPException(status_code=400, detail="YouTube: Video erişilebilir değil veya silinmiş.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"YouTube/Whisper İşleme Hatası: {str(e)}")
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
            temp_dir = os.path.dirname(temp_file_path)
            if os.path.isdir(temp_dir) and not os.listdir(temp_dir):
                os.rmdir(temp_dir)


# --- RAPOR PROMPTU (DEĞİŞİKLİK YOK) ---
RAPOR_PROMPTU = """
(ÖNCEKİ main.py'den RAPOR_PROMPTU metninin tamamını buraya kopyalayın)
"""
RAPOR_PROMPTU = """
Sen, 'Akıllı İçerik Platformu' adına çalışan, detay odaklı bir yapay zekâ uzmanısın. 
Görevin, sana verilen bir içerik metnini analiz etmek ve bu metni, öğrenmeyi ve eyleme geçmeyi kolaylaştıracak şekilde, 8 ana başlıkta özetleyen net bir **Markdown (.md) formatında rapor** hazırlamaktır.
Raporun formatı AŞAĞIDAKİ YAPILANDIRILMIŞ ŞEKİLDE, TÜM BAŞLIKLAR ZORUNLU OLARAK OLMALIDIR:
### 1. Konu Özeti (3-5 Cümle)
[Buraya içeriğin genel amacını ve ana temasını 3-5 net cümle ile yaz.]
### 2. Konu Bölümlendirme (Gezinme Haritası)
[İçeriğin ana başlıklarını ve mantıksal akışını gösteren bir liste hazırla. Bölüm başlıklarını net bir şekilde listele.]
* Bölüm 1 Adı
* Bölüm 2 Adı
* ...
### 3. Temel Kavramlar Sözlüğü (Markdown Tablosu)
["Bu içerikte bilmem gereken kilit kelimeler neler?" sorusuna cevap ver. Bu kavramları ve kısa tanımlarını içeren iki sütunlu bir Markdown Tablosu oluştur. Tabloya en az 5 temel kavram ekle.]
| Kavram | Tanım |
|---|---|
| Bilgi Güvenliği | Hassas verilerin yetkisiz erişime karşı korunması. |
| ... | ... |
### 4. Öğrenme Çıkarımları (Liste)
[Bu içerik bittiğinde aklımda kalması gereken 3 ana prensibi maddeler halinde, kısa ve öz olarak yaz.]
* Ana prensip 1
* Ana prensip 2
* ...
### 5. Pratik Öneri (1 Paragraf)
[Bu bilgiyi gerçek hayatta veya iş akışında nasıl kullanabileceğine dair 1 paragraflık somut bir eylem önerisi yaz.]
### 6. Faydalı Kaynaklar ve Araçlar (Liste)
[İçeriğin konusuyla ilgili, öğrenmeyi derinleştirecek ve işe yarayacak 3-5 adet ek kaynak, araç, program, site veya bu alandaki uzmanların adını bir liste olarak öner.]
* Kaynak/Araç 1 (Kısa açıklama)
* Kaynak/Araç 2 (Kısa açıklama)
* ...
### 7. Mini Quiz (3-5 Soru)
[Kullanıcının konuyu ne kadar anladığını test etmek için 3 adet, kısa cevaplı veya çoktan seçmeli, kritik soru hazırla. Her sorunun hemen altına DOĞRU CEVABI da **parantez içinde** belirt.]
1. Soru 1? (Cevap: ...)
2. Soru 2? (Cevap: ...)
3. Soru 3? (Cevap: ...)
### 8. Kişisel Notlar (Boş Alan)
[Bu bölümü kullanıcı kendi notlarını alsın diye boş bırak. Sadece '### 8. Kişisel Notlar' başlığını yaz ve altını boş bırak.]
"""

# --- YENİ KULLANICI VE OTURUM ENDPOINT'LERİ ---

@app.post("/register", response_model=schemas.TokenResponse, tags=["Kimlik Doğrulama"])
async def register_user(user: schemas.UserCreate, db: Session = Depends(get_db)):
    """
    Yeni kullanıcıyı kalıcı PostgreSQL veritabanına kaydeder
    ve onun için kalıcı bir Token oluşturur.
    """
    # E-posta veya Kullanıcı ID'si zaten var mı diye kontrol et
    db_user = crud.get_user_by_email(db, email=user.email)
    if db_user:
        raise HTTPException(status_code=400, detail="Bu e-posta adresi zaten kayıtlı.")
    
    db_user_by_id = crud.get_user_by_user_id_str(db, user_id_str=user.user_id_str)
    if db_user_by_id:
        raise HTTPException(status_code=400, detail="Bu kullanıcı ID'si zaten kullanılıyor.")
    
    # Yeni kullanıcıyı veritabanında oluştur
    db_user = crud.create_user(db=db, user=user)
    
    # Kullanıcı için yeni bir Token oluştur ve DB'ye kaydet
    new_token_str = secrets.token_urlsafe(32)
    crud.create_user_token(db=db, user=db_user, token=new_token_str)
    
    return {"access_token": new_token_str, "token_type": "bearer"}

@app.post("/token", response_model=schemas.TokenResponse, tags=["Kimlik Doğrulama"])
async def login_for_access_token(
    form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)
):
    """
    Kullanıcının (user_id_str) ve şifresinin (password)
    veritabanındaki hash ile eşleşip eşleşmediğini kontrol eder.
    Eşleşirse, yeni bir token oluşturur.
    NOT: Frontend'in 'Giriş Yap' formu bu endpoint'i kullanmalı.
    """
    # Not: OAuth2PasswordRequestForm "username" alanı bekler.
    # Biz burada "username" alanını "user_id_str" olarak kullanıyoruz.
    user = crud.get_user_by_user_id_str(db, user_id_str=form_data.username)
    
    if not user or not crud.verify_password(form_data.password, user.password_hash):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Geçersiz kullanıcı ID'si veya şifre",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    # Başarılı girişte yeni bir Token oluştur ve DB'ye kaydet
    new_token_str = secrets.token_urlsafe(32)
    crud.create_user_token(db=db, user=user, token=new_token_str)
    
    return {"access_token": new_token_str, "token_type": "bearer"}


@app.get("/users/me", response_model=schemas.User, tags=["Kimlik Doğrulama"])
async def read_users_me(current_user: models.User = Depends(get_current_user_from_token)):
    """
    X-API-TOKEN'ı doğrular ve ilgili kullanıcıyı döndürür.
    Frontend'in "token geçerli mi?" kontrolü (Tespit 3.2) için kullanılır.
    """
    return current_user


# --- YENİ RAPOR ENDPOINT'LERİ (Tespit 3.3 Çözümü) ---

@app.get("/reports/my-reports", response_model=List[schemas.Report], tags=["Raporlama"])
async def read_user_reports(
    current_user: models.User = Depends(get_current_user_from_token),
    db: Session = Depends(get_db)
):
    """
    Geçerli kullanıcının tüm geçmiş raporlarını (GCS linklerini)
    veritabanından listeler.
    """
    reports = crud.get_user_reports(db, user_id=current_user.id)
    return reports


# --- ANA İŞLEM ENDPOINT'i (YENİDEN YAZILDI) ---

@app.post("/analiz-et", tags=["Raporlama"])
async def analiz_et_ve_raporla(
    dosya: Optional[UploadFile] = File(None), 
    youtube_url: Optional[str] = None,       
    current_user: models.User = Depends(get_current_user_from_token),
    db: Session = Depends(get_db) 
):
    
    # 1. GÜVENLİK KONTROLÜ
    # 'current_user' bağımlılığı (Depends) sayesinde token kontrolü zaten yapıldı.
    # 'current_user' objesi elimizde.
    user_id = current_user.id
    user_id_str = current_user.user_id_str

    # 2. İÇERİK TÜRÜNÜ BELİRLEME ve OKUMA (Async Bloklama Çözümü ile)
    metin = ""
    dosya_adi_temel = "Analiz_Raporu"

    try:
        if dosya:
            # Uzantı kontrolü (Dosya boyutu Frontend'de kontrol edildi)
            uzanti = os.path.splitext(dosya.filename)[1].lower()
            if uzanti not in ALLOWED_EXTENSIONS:
                 raise HTTPException(status_code=400, detail="Desteklenmeyen dosya türü.")
            
            dosya_adi_temel = os.path.splitext(dosya.filename)[0]

            # Yavaş I/O işlemlerini threadpool'da çalıştır (Tespit 2.3)
            if uzanti in [".mp3", ".wav", ".m4a"]:
                metin = await run_in_threadpool(read_audio, dosya)
            elif uzanti == ".pdf":
                metin = await run_in_threadpool(read_pdf, dosya)
            elif uzanti in [".docx", ".doc"]:
                metin = await run_in_threadpool(read_docx, dosya)
            elif uzanti in [".pptx", ".ppt"]:
                metin = await run_in_threadpool(read_pptx, dosya)
            elif uzanti in [".jpg", ".jpeg", ".png"]:
                metin = await run_in_threadpool(read_image, dosya)
            else:
                raise HTTPException(status_code=400, detail="Desteklenmeyen dosya türü.")
        
        elif youtube_url:
            metin = await run_in_threadpool(download_youtube_audio, youtube_url)
            dosya_adi_temel = f"youtube-video-analizi"
        
        else:
            raise HTTPException(status_code=400, detail="Dosya yükleyin veya bir YouTube URL'si sağlayın.")
            
    except HTTPException as h:
        raise h 
    except Exception as e:
        print(f"İçerik Okuma Başarısız: {e}")
        raise HTTPException(status_code=500, detail=f"İçerik Okuma Başarısız oldu. {e}")

    # 3. METİN ANALİZİ (GPT-4o) (Async Bloklama Çözümü ile)
    if not metin or metin.strip() == "":
         raise HTTPException(status_code=400, detail="İçerikten metin çıkarılamadı (Dosya boş veya okunamadı).")

    try:
        # OpenAI çağrısı yavaş bir I/O işlemidir, threadpool'a gönder
        def run_openai_call():
            return client.chat.completions.create(
                model="gpt-4o",   
                messages=[
                    {"role": "system", "content": RAPOR_PROMPTU},
                    {"role": "user", "content": f"Lütfen aşağıdaki içerik metnini analiz et ve raporla:\n\n{metin}"}
                ]
            )
        
        chat_completion = await run_in_threadpool(run_openai_call)
        rapor_metni = chat_completion.choices[0].message.content
        print(f"Rapor oluşturuldu. Kullanıcı: {user_id_str}")
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM/Raporlama hatası: {str(e)}")


    # --- 4. KALICI GCS DEPOLAMA (Async Bloklama Çözümü ile) ---
    try:
        # 1. GCS İstemcisini Başlatma (Hızlı)
        sa_key_json = os.getenv(GCS_KEY_ENV_VAR)
        if not sa_key_json:
            raise Exception("GCS Service Account Key çevresel değişkeni ayarlanmadı.")
        credentials_dict = json.loads(sa_key_json)
        gcs_client = storage.Client.from_service_account_info(credentials_dict)
        bucket = gcs_client.bucket(GCS_BUCKET_NAME)

        # 2. Dosya Adını Oluşturma (Hızlı)
        temiz_ad = slugify(dosya_adi_temel)
        zaman_damgasi = datetime.now().strftime("%Y-%m-%d_%H%M%S")
        gcs_file_name = f"{user_id_str}/{temiz_ad}_{zaman_damgasi}.md" # Klasör adı olarak user_id_str kullandık
        
        # 3. GCS'e Yükleme (Yavaş I/O, threadpool'a gönder)
        def run_gcs_upload():
            blob = bucket.blob(gcs_file_name)
            blob.upload_from_string(
                data=rapor_metni.encode('utf-8'), 
                content_type='text/markdown; charset=utf-8'
            )
            print(f"Rapor başarıyla GCS'e yüklendi: {gcs_file_name}")
        
        await run_in_threadpool(run_gcs_upload)
        
        # 4. Genel Erişim URL'sini Oluşturma (Hızlı)
        kaydedilen_dosya_url = f"https://storage.googleapis.com/{GCS_BUCKET_NAME}/{gcs_file_name}"

    except Exception as e:
        print(f"HATA: Rapor GCS'e kaydedilemedi: {e}")
        raise HTTPException(status_code=500, detail=f"Bulut Depolama Hatası: Rapor GCS'e kaydedilemedi. {e}")


    # --- 5. RAPORU VERİTABANINA KAYDETME (Tespit 3.3 Çözümü) ---
    try:
        report_schema = schemas.ReportCreate(
            gcs_url=kaydedilen_dosya_url,
            file_name=gcs_file_name
        )
        crud.create_user_report(db=db, report=report_schema, user_id=user_id)
        print(f"Raporun GCS linki veritabanına kaydedildi. Kullanıcı: {user_id_str}")
    
    except Exception as e:
        # Bu kritik bir hata değil, kullanıcı raporu yine de alır.
        print(f"HATA: Rapor veritabanına kaydedilemedi (ama GCS'e yüklendi): {e}")


    # 6. KULLANICIYA YANIT DÖNÜŞÜ
    return {
        "user_id": user_id_str,
        "rapor_markdown": rapor_metni, 
        "dosya_url": kaydedilen_dosya_url
    }

# --- FRONTEND'İ SUNMA ---
app.mount("/", StaticFiles(directory="frontend", html=True), name="static")
